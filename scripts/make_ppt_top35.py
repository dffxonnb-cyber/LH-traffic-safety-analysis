#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
TOP35 15분 발표용 PowerPoint(.pptx) 생성 스크립트.

- 공개 저장소 버전에서는 발표 가이드 문서를 제외하고 코드만 유지합니다.
- 이미지: data/통합_데이터/top35_outputs/figures/, 시각화_공유 등 (없으면 텍스트만)
- 출력: data/통합_데이터/top35_outputs/LH_TOP35_15min.pptx

실행: 프로젝트 루트(1최종_LH)에서
  python scripts/make_ppt_top35.py
또는
  python scripts/make_ppt_top35.py --out my_presentation.pptx
"""

from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    raise SystemExit("python-pptx 필요: pip install python-pptx")

# 프로젝트 루트 = scripts/ 의 상위
ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data" / "통합_데이터"
TOP35 = DATA / "top35_outputs"
FIGURES = TOP35 / "figures"
VIZ = DATA / "시각화_공유"
NOTES = ROOT / "notebooks"


def _find_image(*candidates: Path) -> Path | None:
    for p in candidates:
        if p is not None and p.exists():
            return p
    return None


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]  # title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    image_path: Path | None = None,
    image_left: float = 0.5,
    image_top: float = 1.8,
    image_width: float = 6.5,
) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)

    # 제목
    left = Inches(0.5)
    top = Inches(0.3)
    w = Inches(9)
    h = Inches(0.7)
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True

    # 본문 불릿
    top = Inches(1.0)
    h = Inches(1.2) if image_path else Inches(5)
    tb2 = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), h)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(bullets):
        if i == 0:
            p2 = tf2.paragraphs[0]
        else:
            p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(14)
        p2.space_after = Pt(6)

    # 이미지(선택)
    if image_path and image_path.exists():
        try:
            slide.shapes.add_picture(
                str(image_path),
                Inches(image_left),
                Inches(image_top),
                width=Inches(image_width),
            )
        except Exception:
            pass


def build_slides(prs: Presentation) -> None:
    # 슬라이드 1: 오프닝
    img1 = _find_image(
        FIGURES / "fig07_transfer_kpi_card.png",
        VIZ / "하남교산_위험점수_지도_분포.png",
    )
    _add_content_slide(
        prs,
        "하남교산, 안전은 완공 후가 아니라 지금 설계해야 합니다",
        [
            "우리는 위험을 예측한 것이 아니라, 설치 순서를 결정했습니다.",
            "하남교산은 아직 완성된 도시가 아닙니다. 도시가 형성되는 지금 위험을 먼저 설계해야 합니다.",
        ],
        img1,
    )

    # 슬라이드 2: 분석 세계관
    img2 = _find_image(VIZ / "4시구_위험점수_지도_분포.png", VIZ / "하남교산_유사고위험_분포.png")
    _add_content_slide(
        prs,
        "4개 시구를 학습해 교산으로 전이",
        [
            "학습(과거)과 설치(미래)를 분리해 의사결정 정확도를 높였습니다.",
            "1) 4개 시구에서 사고 패턴 학습 → 2) 하남교산 우선순위 격자에 전이 → 3) 커버리지 기반 Top20·지점별 시설 패키지 제안",
        ],
        img2,
    )

    # 슬라이드 3: 전이 타당성
    img3 = _find_image(
        FIGURES / "fig01_transfer_by_region.png",
        NOTES / "GWRF_vs_Priority_Correlation.png",
    )
    _add_content_slide(
        prs,
        "지역을 통째로 빼도 성능이 유지된다",
        [
            "평균 AUC 0.8604, Top10 리프트 4.39x, 최저(송파) AUC 0.7979",
            "4개 지역 중 1개를 빼고 학습한 뒤 제외 지역에서 측정. 지역이 달라도 패턴이 유지됨을 확인했습니다.",
        ],
        img3,
    )

    # 슬라이드 4: 원인 설명력
    img4 = _find_image(
        FIGURES / "fig02_feature_stability.png",
        VIZ / "04_5_유형별_우선순위점수.png",
    )
    _add_content_slide(
        prs,
        "교통량·속도·혼잡이 반복적으로 상위",
        [
            "AADT_mean 0.2984, velocity_mean 0.2251, TI_mean 0.1492",
            "반복 학습에서 위 변수들이 상위권 유지. 교통량·속도·혼잡이 위험 형성의 핵심축입니다.",
        ],
        img4,
    )

    # 슬라이드 5: 강건성
    img5 = _find_image(
        FIGURES / "fig03_mc_robustness.png",
        DATA / "hanam_gyosan_selected_sites_k_scenarios.png",
    )
    _add_content_slide(
        prs,
        "가중치가 바뀌어도 추천안은 절반 이상 유지",
        [
            "MC 평균 Jaccard 0.503, 범위 0.429~0.538, 평균 커버리지 0.671",
            "가중치를 바꿔도 결과는 민감하지만 붕괴하지 않습니다.",
        ],
        img5,
    )

    # 슬라이드 6: 정책 선택지
    img6 = _find_image(
        FIGURES / "fig04_scenario_tradeoff.png",
        VIZ / "04_3_설치우선순위_Top20.png",
    )
    _add_content_slide(
        prs,
        "커버리지와 기존안 연속성의 트레이드오프",
        [
            "risk60_flow40: coverage 0.668, Jaccard 0.538",
            "risk70_flow30: coverage 0.671, Jaccard 0.538",
            "risk80_flow20: coverage 0.674, Jaccard 0.481",
            "정책 목표가 연속성이면 60:40/70:30, 확장성이면 80:20.",
        ],
        img6,
    )

    # 슬라이드 7: 공간 설계 결과
    img7 = _find_image(
        DATA / "hanam_gyosan_combined_final_plan.png",
        DATA / "hanam_gyosan_selected_sites_k20.png",
    )
    _add_content_slide(
        prs,
        "취약축 중심 배치로 사각지대 축소",
        [
            "무작위 분산이 아니라 취약축을 따라 배치했습니다.",
            "선정 지점은 취약축과 수요 밀집구간을 따라 배치했고, 중복 커버리지를 줄이면서 사각지대를 메우는 구조입니다.",
        ],
        img7,
    )

    # 슬라이드 8: 시설 패키지
    img8 = _find_image(
        FIGURES / "fig05_top20_confidence.png",
        FIGURES / "fig06_top20_package_mix.png",
    )
    _add_content_slide(
        prs,
        "Top20 지점별 맞춤 패키지",
        [
            "도로 중심 고위험 구간: 교통정온화·스마트 횡단",
            "주거·녹지 구간: CCTV, 비상벨, 조명 보강 조합",
            "gyosan_top20_facility_blueprint.csv: gid, facility, recommended_package, recommendation_reason",
        ],
        img8,
    )

    # 슬라이드 9: 실행 로드맵
    img9 = _find_image(VIZ / "하남교산_우선순위_지도.png", VIZ / "04_6_구역별_위험패턴.png")
    _add_content_slide(
        prs,
        "3단계 집행 프레임",
        [
            "1단계: very_high / high 즉시 검토",
            "2단계: medium 현장 검증 후 반영",
            "3단계: 분기별 재최적화",
            "이번 결과는 단발성 분석이 아니라 운영형 프레임입니다.",
        ],
        img9,
    )

    # 슬라이드 10: 클로징
    img10 = _find_image(
        FIGURES / "fig07_transfer_kpi_card.png",
        VIZ / "하남교산_설치제안.png",
    )
    _add_content_slide(
        prs,
        "전이 타당성 + 강건성 + 실행성",
        [
            "우리는 성능만이 아니라 전이 가능성을 검증했습니다.",
            "우리는 추천만이 아니라 흔들림을 측정했습니다.",
            "우리는 분석만이 아니라 설치 순서를 제안했습니다.",
        ],
        img10,
    )

    # --- Q&A 백업 슬라이드 ---
    _add_content_slide(
        prs,
        "Q&A 백업: 송파 성능이 낮은 이유?",
        [
            "송파는 양성비·공간 구조가 상대적으로 이질적이라 성능이 낮게 나왔습니다.",
            "그럼에도 AUC 0.79대는 전이 적용 관점에서 유효한 범위입니다.",
            "참고: transfer_loro_detail.csv",
        ],
    )
    _add_content_slide(
        prs,
        "Q&A 백업: 가중치 바꾸면 결과 붕괴?",
        [
            "평균 Jaccard 0.503, 하한 0.429로 붕괴 아님.",
            "핵심 후보는 유지됩니다. 참고: gyosan_mc_runs.csv",
        ],
    )
    _add_content_slide(
        prs,
        "Q&A 백업: 왜 여기 CCTV, 저기 비상벨?",
        [
            "지점별 dominant zone + 고위험/노출 조합 기반 패키지.",
            "gyosan_top20_facility_blueprint.csv에서 recommended_package, recommendation_reason으로 추적 가능.",
        ],
    )
    _add_content_slide(
        prs,
        "Q&A 백업: 한계와 보완",
        [
            "시야·전원·민원·공사일정 등 현장 제약 검증 필수.",
            "운영데이터 갱신 시 우선순위 일부 이동 가능 → 분기 단위 재최적화를 운영 프로세스에 포함.",
        ],
    )


def main() -> None:
    parser = argparse.ArgumentParser(description="TOP35 15분 발표용 pptx 생성")
    parser.add_argument(
        "--out",
        type=Path,
        default=TOP35 / "LH_TOP35_15min.pptx",
        help="출력 .pptx 경로",
    )
    args = parser.parse_args()

    out_path = args.out
    if not out_path.is_absolute():
        out_path = ROOT / out_path
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    build_slides(prs)
    prs.save(str(out_path))
    print(f"[DONE] 저장: {out_path}")
    print("  슬라이드: 본문 10장 + Q&A 백업 4장")


if __name__ == "__main__":
    main()
